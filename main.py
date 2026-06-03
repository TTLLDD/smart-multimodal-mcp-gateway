from __future__ import annotations

import logging
import os
import re
import time
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

import httpx
import uvicorn
from dotenv import load_dotenv
from fastapi import FastAPI, HTTPException
from fastapi.responses import HTMLResponse
from pydantic import BaseModel, Field, field_validator
from mcp.server.fastmcp import FastMCP


load_dotenv()

logger = logging.getLogger("smart_multimodal_mcp_gateway")
logging.basicConfig(level=logging.INFO, format="%(asctime)s %(levelname)s %(message)s")

ENV_PATH = Path(__file__).with_name(".env")
MODELS_PATH = Path(__file__).with_name("models.json")
OUTPUTS_PATH = Path(__file__).with_name("outputs")
SKILLS_PATH = Path(__file__).with_name("skills")
KEEP_SECRET = "__KEEP_SECRET__"

CONFIG_FIELDS: Dict[str, str] = {
    "MCP_SERVER_NAME": "smart-multimodal-mcp-gateway",
    "HOST": "0.0.0.0",
    "PORT": "8010",
    "DEEPSEEK_API_KEY": "",
    "DEEPSEEK_BASE_URL": "https://api.deepseek.com/v1",
    "DEEPSEEK_TEXT_MODEL": "deepseek-chat",
    "DASHSCOPE_API_KEY": "",
    "DASHSCOPE_BASE_URL": "https://dashscope.aliyuncs.com/compatible-mode/v1",
    "DASHSCOPE_VISION_MODEL": "qwen-vl-max",
    "OPENAI_API_KEY": "",
    "OPENAI_BASE_URL": "https://api.openai.com/v1",
    "OPENAI_VISION_MODEL": "gpt-4o",
    "DEFAULT_TEMPERATURE": "0.7",
    "DEFAULT_TOP_P": "1.0",
    "DEFAULT_MAX_TOKENS": "2048",
    "HTTP_CONNECT_TIMEOUT": "10",
    "HTTP_READ_TIMEOUT": "120",
    "HTTP_WRITE_TIMEOUT": "120",
    "HTTP_POOL_TIMEOUT": "30",
    "HTTP_MAX_CONNECTIONS": "100",
    "HTTP_MAX_KEEPALIVE_CONNECTIONS": "20",
}


def env_str(name: str, default: str = "") -> str:
    return os.getenv(name, default).strip()


def env_int(name: str, default: int) -> int:
    raw = os.getenv(name)
    if not raw:
        return default
    try:
        return int(raw)
    except ValueError:
        return default


def env_float(name: str, default: float) -> float:
    raw = os.getenv(name)
    if not raw:
        return default
    try:
        return float(raw)
    except ValueError:
        return default


class GatewayConfig(BaseModel):
    MCP_SERVER_NAME: str = Field(default=CONFIG_FIELDS["MCP_SERVER_NAME"])
    HOST: str = Field(default=CONFIG_FIELDS["HOST"])
    PORT: str = Field(default=CONFIG_FIELDS["PORT"])
    DEEPSEEK_API_KEY: str = ""
    DEEPSEEK_BASE_URL: str = Field(default=CONFIG_FIELDS["DEEPSEEK_BASE_URL"])
    DEEPSEEK_TEXT_MODEL: str = Field(default=CONFIG_FIELDS["DEEPSEEK_TEXT_MODEL"])
    DASHSCOPE_API_KEY: str = ""
    DASHSCOPE_BASE_URL: str = Field(default=CONFIG_FIELDS["DASHSCOPE_BASE_URL"])
    DASHSCOPE_VISION_MODEL: str = Field(default=CONFIG_FIELDS["DASHSCOPE_VISION_MODEL"])
    OPENAI_API_KEY: str = ""
    OPENAI_BASE_URL: str = Field(default=CONFIG_FIELDS["OPENAI_BASE_URL"])
    OPENAI_VISION_MODEL: str = Field(default=CONFIG_FIELDS["OPENAI_VISION_MODEL"])
    DEFAULT_TEMPERATURE: str = Field(default=CONFIG_FIELDS["DEFAULT_TEMPERATURE"])
    DEFAULT_TOP_P: str = Field(default=CONFIG_FIELDS["DEFAULT_TOP_P"])
    DEFAULT_MAX_TOKENS: str = Field(default=CONFIG_FIELDS["DEFAULT_MAX_TOKENS"])
    HTTP_CONNECT_TIMEOUT: str = Field(default=CONFIG_FIELDS["HTTP_CONNECT_TIMEOUT"])
    HTTP_READ_TIMEOUT: str = Field(default=CONFIG_FIELDS["HTTP_READ_TIMEOUT"])
    HTTP_WRITE_TIMEOUT: str = Field(default=CONFIG_FIELDS["HTTP_WRITE_TIMEOUT"])
    HTTP_POOL_TIMEOUT: str = Field(default=CONFIG_FIELDS["HTTP_POOL_TIMEOUT"])
    HTTP_MAX_CONNECTIONS: str = Field(default=CONFIG_FIELDS["HTTP_MAX_CONNECTIONS"])
    HTTP_MAX_KEEPALIVE_CONNECTIONS: str = Field(
        default=CONFIG_FIELDS["HTTP_MAX_KEEPALIVE_CONNECTIONS"]
    )


class ModelProfile(BaseModel):
    id: str = ""
    display_name: str = ""
    capability: str = "text"
    base_url: str = ""
    api_key: str = ""
    model: str = ""
    model_options: List[str] = Field(default_factory=list)
    enabled: bool = True

    @field_validator("capability")
    @classmethod
    def validate_capability(cls, value: str) -> str:
        normalized = (value or "").strip().lower()
        if normalized not in {"text", "multimodal", "fallback"}:
            raise ValueError("capability must be text, multimodal, or fallback")
        return normalized


class ModelSettings(BaseModel):
    profiles: List[ModelProfile] = Field(default_factory=list)
    default_text_profile_id: str = ""
    default_multimodal_profile_id: str = ""


class ModelTestRequest(BaseModel):
    profile: Optional[ModelProfile] = None
    profile_id: Optional[str] = None


def parse_env_file(path: Path = ENV_PATH) -> Dict[str, str]:
    values: Dict[str, str] = {}
    if not path.exists():
        return values

    for line in path.read_text(encoding="utf-8").splitlines():
        stripped = line.strip()
        if not stripped or stripped.startswith("#") or "=" not in stripped:
            continue
        key, value = stripped.split("=", 1)
        key = key.strip()
        value = value.strip()
        if len(value) >= 2 and value[0] == value[-1] and value[0] in {"'", '"'}:
            value = value[1:-1]
        values[key] = value
    return values


def merged_config_values() -> Dict[str, str]:
    file_values = parse_env_file()
    merged = dict(CONFIG_FIELDS)
    merged.update(file_values)
    for key in CONFIG_FIELDS:
        if key in os.environ:
            merged[key] = os.environ[key]
    return merged


def quote_env_value(value: str) -> str:
    value = "" if value is None else str(value)
    if not value:
        return ""
    if any(char in value for char in (" ", "\t", "\n", "\r", '"', "'", "#")):
        return '"' + value.replace("\\", "\\\\").replace('"', '\\"') + '"'
    return value


def write_env_file(config: GatewayConfig) -> None:
    values = config.model_dump()
    lines = [
        "# Smart Multimodal MCP SSE Gateway",
        "# Local BYOK config. Do not commit real API keys.",
        "",
        "MCP_SERVER_NAME=" + quote_env_value(values["MCP_SERVER_NAME"]),
        "HOST=" + quote_env_value(values["HOST"]),
        "PORT=" + quote_env_value(values["PORT"]),
        "",
        "# Text-first provider: DeepSeek",
        "DEEPSEEK_API_KEY=" + quote_env_value(values["DEEPSEEK_API_KEY"]),
        "DEEPSEEK_BASE_URL=" + quote_env_value(values["DEEPSEEK_BASE_URL"]),
        "DEEPSEEK_TEXT_MODEL=" + quote_env_value(values["DEEPSEEK_TEXT_MODEL"]),
        "",
        "# Multimodal-first provider: Alibaba DashScope OpenAI-compatible mode",
        "DASHSCOPE_API_KEY=" + quote_env_value(values["DASHSCOPE_API_KEY"]),
        "DASHSCOPE_BASE_URL=" + quote_env_value(values["DASHSCOPE_BASE_URL"]),
        "DASHSCOPE_VISION_MODEL=" + quote_env_value(values["DASHSCOPE_VISION_MODEL"]),
        "",
        "# Optional multimodal provider: OpenAI or a compatible relay",
        "OPENAI_API_KEY=" + quote_env_value(values["OPENAI_API_KEY"]),
        "OPENAI_BASE_URL=" + quote_env_value(values["OPENAI_BASE_URL"]),
        "OPENAI_VISION_MODEL=" + quote_env_value(values["OPENAI_VISION_MODEL"]),
        "",
        "# Generation defaults",
        "DEFAULT_TEMPERATURE=" + quote_env_value(values["DEFAULT_TEMPERATURE"]),
        "DEFAULT_TOP_P=" + quote_env_value(values["DEFAULT_TOP_P"]),
        "DEFAULT_MAX_TOKENS=" + quote_env_value(values["DEFAULT_MAX_TOKENS"]),
        "",
        "# HTTPX settings for large multimodal payloads",
        "HTTP_CONNECT_TIMEOUT=" + quote_env_value(values["HTTP_CONNECT_TIMEOUT"]),
        "HTTP_READ_TIMEOUT=" + quote_env_value(values["HTTP_READ_TIMEOUT"]),
        "HTTP_WRITE_TIMEOUT=" + quote_env_value(values["HTTP_WRITE_TIMEOUT"]),
        "HTTP_POOL_TIMEOUT=" + quote_env_value(values["HTTP_POOL_TIMEOUT"]),
        "HTTP_MAX_CONNECTIONS=" + quote_env_value(values["HTTP_MAX_CONNECTIONS"]),
        "HTTP_MAX_KEEPALIVE_CONNECTIONS="
        + quote_env_value(values["HTTP_MAX_KEEPALIVE_CONNECTIONS"]),
        "",
    ]
    ENV_PATH.write_text("\n".join(lines), encoding="utf-8")


def merge_secret_placeholders(config: GatewayConfig) -> GatewayConfig:
    current = merged_config_values()
    values = config.model_dump()
    for key in ("DEEPSEEK_API_KEY", "DASHSCOPE_API_KEY", "OPENAI_API_KEY"):
        submitted = values.get(key, "")
        if submitted == "__KEEP_SECRET__" or submitted == "":
            values[key] = current.get(key, "")
    return GatewayConfig(**values)


def apply_config_to_process(config: GatewayConfig) -> None:
    for key, value in config.model_dump().items():
        os.environ[key] = str(value)


def masked_values(values: Dict[str, str]) -> Dict[str, str]:
    result = dict(values)
    for key in ("DEEPSEEK_API_KEY", "DASHSCOPE_API_KEY", "OPENAI_API_KEY"):
        if result.get(key):
            result[key] = "__KEEP_SECRET__"
    return result


def public_config() -> Dict[str, Any]:
    values = merged_config_values()
    return {
        "values": masked_values(values),
        "env_path": str(ENV_PATH),
        "env_exists": ENV_PATH.exists(),
        "configured": {
            "deepseek": bool(values.get("DEEPSEEK_API_KEY")),
            "dashscope": bool(values.get("DASHSCOPE_API_KEY")),
            "openai": bool(values.get("OPENAI_API_KEY")),
        },
        "restart_note": (
            "Saved settings are applied to provider routing immediately. "
            "HOST, PORT, and MCP_SERVER_NAME require a service restart."
        ),
    }


def json_loads(text: str) -> Any:
    import json

    return json.loads(text)


def json_dumps(value: Any) -> str:
    import json

    return json.dumps(value, ensure_ascii=False, indent=2)


def extract_json_object(text: str) -> Dict[str, Any]:
    import json

    raw = (text or "").strip()
    if raw.startswith("```"):
        raw = re.sub(r"^```(?:json)?\s*", "", raw, flags=re.IGNORECASE)
        raw = re.sub(r"\s*```$", "", raw)

    try:
        data = json.loads(raw)
    except json.JSONDecodeError:
        start = raw.find("{")
        end = raw.rfind("}")
        if start < 0 or end <= start:
            raise ValueError("Model did not return a JSON object.")
        data = json.loads(raw[start : end + 1])

    if not isinstance(data, dict):
        raise ValueError("Model JSON must be an object.")
    return data


def safe_filename(value: str, suffix: str) -> str:
    name = re.sub(r"[^\w\u4e00-\u9fff.-]+", "_", (value or "").strip(), flags=re.UNICODE)
    name = name.strip("._") or "output"
    if not name.lower().endswith(suffix.lower()):
        name += suffix
    return name[:120]


def normalize_slide_items(value: Any) -> List[str]:
    if isinstance(value, list):
        return [str(item).strip() for item in value if str(item).strip()][:6]
    if isinstance(value, str):
        return [line.strip(" -•\t") for line in value.splitlines() if line.strip()][:6]
    return []


def build_pptx_file(title: str, slides: List[Dict[str, Any]], output_name: Optional[str] = None) -> Path:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    OUTPUTS_PATH.mkdir(parents=True, exist_ok=True)
    path = OUTPUTS_PATH / safe_filename(output_name or title or "presentation", ".pptx")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    accent = RGBColor(20, 118, 212)
    dark = RGBColor(16, 42, 67)
    muted = RGBColor(95, 113, 137)
    pale = RGBColor(243, 248, 255)

    def add_footer(slide: Any, index: int) -> None:
        bar = slide.shapes.add_shape(1, Inches(0), Inches(7.18), Inches(13.333), Inches(0.32))
        bar.fill.solid()
        bar.fill.fore_color.rgb = pale
        bar.line.color.rgb = pale
        page = slide.shapes.add_textbox(Inches(12.15), Inches(7.2), Inches(0.9), Inches(0.18))
        frame = page.text_frame
        frame.text = str(index)
        paragraph = frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.RIGHT
        paragraph.runs[0].font.size = Pt(9)
        paragraph.runs[0].font.color.rgb = muted

    cover = prs.slides.add_slide(blank)
    stripe = cover.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.24), Inches(7.5))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent
    stripe.line.color.rgb = accent
    cover_title = cover.shapes.add_textbox(Inches(0.9), Inches(2.35), Inches(11.2), Inches(1.0))
    cover_title.text_frame.text = title or "Presentation"
    run = cover_title.text_frame.paragraphs[0].runs[0]
    run.font.size = Pt(38)
    run.font.bold = True
    run.font.color.rgb = dark
    subtitle = cover.shapes.add_textbox(Inches(0.94), Inches(3.45), Inches(10.5), Inches(0.55))
    subtitle.text_frame.text = "Generated by local MCP gateway"
    subtitle_run = subtitle.text_frame.paragraphs[0].runs[0]
    subtitle_run.font.size = Pt(16)
    subtitle_run.font.color.rgb = muted
    add_footer(cover, 1)

    for index, slide_data in enumerate(slides, start=2):
        slide = prs.slides.add_slide(blank)
        header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.18))
        header.fill.solid()
        header.fill.fore_color.rgb = accent
        header.line.color.rgb = accent

        heading = str(slide_data.get("title") or f"Slide {index - 1}").strip()
        title_box = slide.shapes.add_textbox(Inches(0.74), Inches(0.62), Inches(11.8), Inches(0.55))
        title_box.text_frame.text = heading
        title_run = title_box.text_frame.paragraphs[0].runs[0]
        title_run.font.size = Pt(27)
        title_run.font.bold = True
        title_run.font.color.rgb = dark

        bullets = normalize_slide_items(slide_data.get("bullets"))
        body = slide.shapes.add_textbox(Inches(0.98), Inches(1.55), Inches(11.2), Inches(4.75))
        frame = body.text_frame
        frame.clear()
        for bullet_index, bullet in enumerate(bullets or [" "]):
            paragraph = frame.paragraphs[0] if bullet_index == 0 else frame.add_paragraph()
            paragraph.text = bullet
            paragraph.level = 0
            paragraph.space_after = Pt(10)
            paragraph.font.size = Pt(20)
            paragraph.font.color.rgb = dark

        note = str(slide_data.get("speaker_note") or slide_data.get("note") or "").strip()
        if note:
            note_box = slide.shapes.add_textbox(Inches(0.98), Inches(6.35), Inches(11.0), Inches(0.42))
            note_box.text_frame.text = note[:180]
            note_run = note_box.text_frame.paragraphs[0].runs[0]
            note_run.font.size = Pt(11)
            note_run.font.color.rgb = muted
        add_footer(slide, index)

    prs.save(path)
    return path


def skill_file_path(skill_name: str) -> Path:
    normalized = slugify(skill_name).replace("-", "_")
    path = SKILLS_PATH / normalized / "SKILL.md"
    if not path.exists():
        raise ValueError(f"Skill not found: {skill_name}")
    return path


def read_skill(skill_name: str) -> Dict[str, str]:
    path = skill_file_path(skill_name)
    text = path.read_text(encoding="utf-8")
    first_line = next((line.strip() for line in text.splitlines() if line.strip()), "")
    title = first_line.lstrip("# ").strip() if first_line.startswith("#") else path.parent.name
    return {
        "name": path.parent.name,
        "title": title,
        "path": str(path),
        "instructions": text,
    }


def list_local_skills() -> List[Dict[str, str]]:
    if not SKILLS_PATH.exists():
        return []
    skills: List[Dict[str, str]] = []
    for skill_dir in sorted(SKILLS_PATH.iterdir()):
        skill_path = skill_dir / "SKILL.md"
        if not skill_path.exists():
            continue
        try:
            skill = read_skill(skill_dir.name)
            summary = ""
            for line in skill["instructions"].splitlines():
                if line.lower().startswith("summary:"):
                    summary = line.split(":", 1)[1].strip()
                    break
            skills.append({
                "name": skill["name"],
                "title": skill["title"],
                "summary": summary,
            })
        except Exception:
            continue
    return skills


def slugify(value: str) -> str:
    slug = "".join(ch.lower() if ch.isalnum() else "-" for ch in value.strip())
    slug = "-".join(part for part in slug.split("-") if part)
    return slug or "model"


def unique_profile_id(base: str, used: set[str]) -> str:
    root = slugify(base)
    candidate = root
    counter = 2
    while candidate in used:
        candidate = f"{root}-{counter}"
        counter += 1
    used.add(candidate)
    return candidate


def legacy_model_settings() -> ModelSettings:
    values = merged_config_values()
    profiles = [
        ModelProfile(
            id="text-default",
            display_name="文本模型",
            capability="text",
            base_url=values["DEEPSEEK_BASE_URL"],
            api_key=values["DEEPSEEK_API_KEY"],
            model=values["DEEPSEEK_TEXT_MODEL"],
            model_options=[values["DEEPSEEK_TEXT_MODEL"]],
            enabled=bool(values["DEEPSEEK_BASE_URL"] and values["DEEPSEEK_TEXT_MODEL"]),
        ),
        ModelProfile(
            id="multimodal-default",
            display_name="多模态模型",
            capability="multimodal",
            base_url=values["DASHSCOPE_BASE_URL"],
            api_key=values["DASHSCOPE_API_KEY"],
            model=values["DASHSCOPE_VISION_MODEL"],
            model_options=[values["DASHSCOPE_VISION_MODEL"]],
            enabled=bool(values["DASHSCOPE_BASE_URL"] and values["DASHSCOPE_VISION_MODEL"]),
        ),
        ModelProfile(
            id="fallback-default",
            display_name="备用多模态",
            capability="fallback",
            base_url=values["OPENAI_BASE_URL"],
            api_key=values["OPENAI_API_KEY"],
            model=values["OPENAI_VISION_MODEL"],
            model_options=[values["OPENAI_VISION_MODEL"]],
            enabled=bool(values["OPENAI_BASE_URL"] and values["OPENAI_VISION_MODEL"]),
        ),
    ]
    return ModelSettings(
        profiles=profiles,
        default_text_profile_id="text-default",
        default_multimodal_profile_id="multimodal-default",
    )


def mask_model_settings(settings: ModelSettings) -> ModelSettings:
    profiles = []
    for profile in settings.profiles:
        data = profile.model_dump()
        if data.get("api_key"):
            data["api_key"] = KEEP_SECRET
        profiles.append(ModelProfile(**data))
    return ModelSettings(
        profiles=profiles,
        default_text_profile_id=settings.default_text_profile_id,
        default_multimodal_profile_id=settings.default_multimodal_profile_id,
    )


def read_model_settings(include_secrets: bool = True) -> ModelSettings:
    if not MODELS_PATH.exists():
        settings = legacy_model_settings()
        return settings if include_secrets else mask_model_settings(settings)

    try:
        data = json_loads(MODELS_PATH.read_text(encoding="utf-8"))
        settings = normalize_model_settings(ModelSettings(**data))
    except Exception as exc:
        logger.warning("Failed to read models.json, using legacy config: %s", exc)
        settings = legacy_model_settings()

    return settings if include_secrets else mask_model_settings(settings)


def normalize_model_settings(settings: ModelSettings, old: Optional[ModelSettings] = None) -> ModelSettings:
    old_by_id = {profile.id: profile for profile in (old.profiles if old else [])}
    used: set[str] = set()
    profiles: List[ModelProfile] = []

    for profile in settings.profiles:
        data = profile.model_dump()
        if not data["id"]:
            data["id"] = unique_profile_id(data["display_name"] or data["model"], used)
        elif data["id"] in used:
            data["id"] = unique_profile_id(data["id"], used)
        else:
            used.add(data["id"])

        old_profile = old_by_id.get(data["id"])
        if data["api_key"] in {"", KEEP_SECRET} and old_profile:
            data["api_key"] = old_profile.api_key

        data["display_name"] = data["display_name"].strip() or data["model"].strip() or data["id"]
        data["base_url"] = data["base_url"].strip()
        data["model"] = data["model"].strip()
        cleaned_options = []
        for option in data.get("model_options", []):
            option = str(option).strip()
            if option and option not in cleaned_options:
                cleaned_options.append(option)
        if data["model"] and data["model"] not in cleaned_options:
            cleaned_options.insert(0, data["model"])
        data["model_options"] = cleaned_options
        profiles.append(ModelProfile(**data))

    profile_ids = {profile.id for profile in profiles}
    enabled = [profile for profile in profiles if profile.enabled]
    text_profiles = [profile for profile in enabled if profile.capability == "text"]
    multimodal_profiles = [profile for profile in enabled if profile.capability == "multimodal"]
    fallback_profiles = [profile for profile in enabled if profile.capability == "fallback"]

    default_text = settings.default_text_profile_id
    if default_text not in profile_ids:
        default_text = (text_profiles or enabled or profiles or [ModelProfile(id="")])[0].id

    default_multi = settings.default_multimodal_profile_id
    if default_multi not in profile_ids:
        default_multi = (multimodal_profiles or fallback_profiles or enabled or profiles or [ModelProfile(id="")])[0].id

    return ModelSettings(
        profiles=profiles,
        default_text_profile_id=default_text,
        default_multimodal_profile_id=default_multi,
    )


def write_model_settings(settings: ModelSettings) -> ModelSettings:
    old = read_model_settings(include_secrets=True)
    normalized = normalize_model_settings(settings, old)
    MODELS_PATH.write_text(json_dumps(normalized.model_dump()), encoding="utf-8")
    return normalized


def public_model_settings() -> Dict[str, Any]:
    secret_settings = read_model_settings(include_secrets=True)
    return {
        "models_path": str(MODELS_PATH),
        "models_exists": MODELS_PATH.exists(),
        "settings": mask_model_settings(secret_settings).model_dump(),
        "configured": {
            profile.id: bool(profile.api_key)
            for profile in secret_settings.profiles
        },
    }


@dataclass(frozen=True)
class ProviderConfig:
    name: str
    base_url: str
    api_key: str
    default_model: str

    @property
    def chat_completions_url(self) -> str:
        return f"{self.base_url.rstrip('/')}/chat/completions"


SERVER_NAME = env_str("MCP_SERVER_NAME", "smart-multimodal-mcp-gateway")

DEFAULT_TEMPERATURE = env_float("DEFAULT_TEMPERATURE", 0.7)
DEFAULT_TOP_P = env_float("DEFAULT_TOP_P", 1.0)
DEFAULT_MAX_TOKENS = env_int("DEFAULT_MAX_TOKENS", 2048)

HTTP_TIMEOUT = httpx.Timeout(
    connect=env_float("HTTP_CONNECT_TIMEOUT", 10.0),
    read=env_float("HTTP_READ_TIMEOUT", 120.0),
    write=env_float("HTTP_WRITE_TIMEOUT", 120.0),
    pool=env_float("HTTP_POOL_TIMEOUT", 30.0),
)
HTTP_LIMITS = httpx.Limits(
    max_connections=env_int("HTTP_MAX_CONNECTIONS", 100),
    max_keepalive_connections=env_int("HTTP_MAX_KEEPALIVE_CONNECTIONS", 20),
)


def provider_config(provider_name: str) -> ProviderConfig:
    provider = provider_name.lower()
    if provider == "deepseek":
        return ProviderConfig(
            name="deepseek",
            base_url=env_str("DEEPSEEK_BASE_URL", "https://api.deepseek.com/v1"),
            api_key=env_str("DEEPSEEK_API_KEY"),
            default_model=env_str("DEEPSEEK_TEXT_MODEL", "deepseek-chat"),
        )
    if provider == "dashscope":
        return ProviderConfig(
            name="dashscope",
            base_url=env_str(
                "DASHSCOPE_BASE_URL",
                "https://dashscope.aliyuncs.com/compatible-mode/v1",
            ),
            api_key=env_str("DASHSCOPE_API_KEY"),
            default_model=env_str("DASHSCOPE_VISION_MODEL", "qwen-vl-max"),
        )
    if provider == "openai":
        return ProviderConfig(
            name="openai",
            base_url=env_str("OPENAI_BASE_URL", "https://api.openai.com/v1"),
            api_key=env_str("OPENAI_API_KEY"),
            default_model=env_str("OPENAI_VISION_MODEL", "gpt-4o"),
        )
    raise ValueError(f"Unsupported provider: {provider_name}")


def provider_from_profile(profile: ModelProfile) -> ProviderConfig:
    return ProviderConfig(
        name=profile.id or profile.display_name or "model-profile",
        base_url=profile.base_url,
        api_key=profile.api_key,
        default_model=profile.model,
    )


def enabled_profiles(settings: ModelSettings) -> List[ModelProfile]:
    return [
        profile
        for profile in settings.profiles
        if profile.enabled and profile.base_url and profile.model
    ]


def find_profile_by_model(settings: ModelSettings, requested_model: Optional[str]) -> Optional[ModelProfile]:
    model_hint = (requested_model or "").strip().lower()
    if not model_hint:
        return None

    for profile in enabled_profiles(settings):
        aliases = {profile.id.lower(), profile.display_name.lower(), profile.model.lower()}
        aliases.update(option.lower() for option in profile.model_options)
        if model_hint in aliases:
            return profile
    return None


def find_profile_by_id(settings: ModelSettings, profile_id: str) -> Optional[ModelProfile]:
    for profile in settings.profiles:
        if profile.id == profile_id:
            return profile
    return None


def default_profile(settings: ModelSettings, has_images: bool) -> Optional[ModelProfile]:
    preferred_id = settings.default_multimodal_profile_id if has_images else settings.default_text_profile_id
    preferred = find_profile_by_id(settings, preferred_id)
    if preferred and preferred.enabled and preferred.base_url and preferred.model:
        return preferred

    profiles = enabled_profiles(settings)
    if has_images:
        for capability in ("multimodal", "fallback", "text"):
            match = next((profile for profile in profiles if profile.capability == capability), None)
            if match:
                return match
    else:
        for capability in ("text", "fallback", "multimodal"):
            match = next((profile for profile in profiles if profile.capability == capability), None)
            if match:
                return match
    return None


def is_image_part(part: Any) -> bool:
    if not isinstance(part, dict):
        return False

    part_type = str(part.get("type", "")).lower()
    if part_type in {"image_url", "input_image"}:
        return True

    image_url = part.get("image_url")
    if isinstance(image_url, dict) and isinstance(image_url.get("url"), str):
        return True

    if isinstance(part.get("url"), str):
        url = part["url"].lower()
        return url.startswith("data:image/") or url.startswith("http://") or url.startswith("https://")

    return False


def content_has_image(content: Any) -> bool:
    if isinstance(content, list):
        return any(is_image_part(part) for part in content)
    if isinstance(content, dict):
        return is_image_part(content)
    if isinstance(content, str):
        return content.strip().lower().startswith("data:image/")
    return False


def messages_have_images(messages: List[Dict[str, Any]]) -> bool:
    return any(content_has_image(message.get("content")) for message in messages if isinstance(message, dict))


def image_to_part(image: str) -> Dict[str, Any]:
    if not isinstance(image, str) or not image.strip():
        raise ValueError("Each image must be a non-empty URL or data:image/...;base64,... string.")
    return {"type": "image_url", "image_url": {"url": image}}


def build_user_content(prompt: Optional[str], images: List[str]) -> Any:
    has_prompt = isinstance(prompt, str) and bool(prompt.strip())
    if images:
        parts: List[Dict[str, Any]] = []
        if has_prompt:
            parts.append({"type": "text", "text": prompt})
        parts.extend(image_to_part(image) for image in images)
        return parts
    if has_prompt:
        return prompt
    raise ValueError("prompt or messages is required.")


def build_messages(
    prompt: Optional[str],
    messages: Optional[List[Dict[str, Any]]],
    images: Optional[List[str]],
) -> List[Dict[str, Any]]:
    image_list = images or []

    if messages is not None and not isinstance(messages, list):
        raise ValueError("messages must be a list of OpenAI-compatible chat messages.")
    if images is not None and not isinstance(images, list):
        raise ValueError("images must be a list of image URLs or data:image/...;base64,... strings.")

    if messages:
        for message in messages:
            if not isinstance(message, dict):
                raise ValueError("Each item in messages must be an object.")
        built = list(messages)
        if prompt or image_list:
            built.append({"role": "user", "content": build_user_content(prompt, image_list)})
        return built

    if prompt or image_list:
        return [{"role": "user", "content": build_user_content(prompt, image_list)}]

    raise ValueError("prompt or messages is required.")


def route_provider(requested_model: Optional[str], has_images: bool) -> Tuple[ProviderConfig, str, str]:
    settings = read_model_settings(include_secrets=True)
    matched_profile = find_profile_by_model(settings, requested_model)
    if matched_profile:
        provider = provider_from_profile(matched_profile)
        return provider, matched_profile.model, f"matched model profile: {matched_profile.display_name}"

    selected_profile = default_profile(settings, has_images)
    if selected_profile:
        provider = provider_from_profile(selected_profile)
        default_kind = "multimodal" if has_images else "text"
        return provider, selected_profile.model, f"default {default_kind} model profile: {selected_profile.display_name}"

    model_hint = (requested_model or "").strip()
    model_lower = model_hint.lower()

    if has_images:
        if any(token in model_lower for token in ("gpt-4o", "gpt-4.1", "openai")):
            provider = provider_config("openai")
            return provider, model_hint or provider.default_model, "image input with OpenAI vision model hint"

        provider = provider_config("dashscope")
        if "qwen" in model_lower or "vl" in model_lower:
            routed_model = model_hint
        else:
            routed_model = provider.default_model
        return provider, routed_model, "image input routes to DashScope vision by default"

    if "deepseek" in model_lower:
        provider = provider_config("deepseek")
        return provider, model_hint or provider.default_model, "DeepSeek model hint"

    if "qwen" in model_lower or "vl" in model_lower:
        provider = provider_config("dashscope")
        return provider, model_hint or provider.default_model, "Qwen/VL model hint"

    if any(token in model_lower for token in ("gpt-4o", "gpt-4.1", "openai")):
        provider = provider_config("openai")
        return provider, model_hint or provider.default_model, "OpenAI model hint"

    provider = provider_config("deepseek")
    return provider, provider.default_model, "text-only default route"


def apply_payload_defaults(provider: ProviderConfig, payload: Dict[str, Any]) -> None:
    model = str(payload.get("model") or "").lower()
    base_url = provider.base_url.lower()
    is_glm = model.startswith("glm-") or "bigmodel.cn" in base_url or "z.ai" in base_url
    if is_glm and "thinking" not in payload:
        # GLM-5.1 enables thinking by default; tool calls need final content.
        payload["thinking"] = {"type": "disabled"}


def extract_text_from_response(data: Dict[str, Any]) -> str:
    choices = data.get("choices")
    if not isinstance(choices, list) or not choices:
        raise RuntimeError("Upstream response did not contain choices.")

    first_choice = choices[0]
    if not isinstance(first_choice, dict):
        raise RuntimeError("Upstream response choice was malformed.")

    message = first_choice.get("message")
    if isinstance(message, dict):
        content = message.get("content")
        if isinstance(content, str):
            return content
        if isinstance(content, list):
            text_parts: List[str] = []
            for part in content:
                if isinstance(part, dict) and isinstance(part.get("text"), str):
                    text_parts.append(part["text"])
                elif isinstance(part, str):
                    text_parts.append(part)
            return "\n".join(text_parts)

    text = first_choice.get("text")
    if isinstance(text, str):
        return text

    raise RuntimeError("Upstream response did not contain text content.")


def brief_upstream_error(response: httpx.Response) -> str:
    try:
        data = response.json()
    except ValueError:
        text = response.text.strip()
        return text[:1000] if text else "No response body."

    if isinstance(data, dict):
        error = data.get("error")
        if isinstance(error, dict):
            message = error.get("message") or error.get("code") or error
            return str(message)[:1000]
        return str(data)[:1000]
    return str(data)[:1000]


_http_client: Optional[httpx.AsyncClient] = None


async def get_http_client() -> httpx.AsyncClient:
    global _http_client
    if _http_client is None or _http_client.is_closed:
        _http_client = httpx.AsyncClient(timeout=HTTP_TIMEOUT, limits=HTTP_LIMITS)
    return _http_client


async def close_http_client() -> None:
    global _http_client
    if _http_client is not None and not _http_client.is_closed:
        await _http_client.aclose()
    _http_client = None


async def call_upstream(provider: ProviderConfig, payload: Dict[str, Any]) -> Dict[str, Any]:
    if not provider.api_key:
        raise RuntimeError(
            f"{provider.name} API key is missing. Please configure it in .env before using this route."
        )

    apply_payload_defaults(provider, payload)
    client = await get_http_client()
    headers = {
        "Authorization": f"Bearer {provider.api_key}",
        "Content-Type": "application/json",
        "Accept": "application/json",
    }

    try:
        response = await client.post(provider.chat_completions_url, headers=headers, json=payload)
    except httpx.TimeoutException as exc:
        raise RuntimeError(f"{provider.name} request timed out: {exc}") from exc
    except httpx.HTTPError as exc:
        raise RuntimeError(f"{provider.name} request failed: {exc}") from exc

    if response.status_code >= 400:
        raise RuntimeError(
            f"{provider.name} upstream error {response.status_code}: {brief_upstream_error(response)}"
        )

    try:
        data = response.json()
    except ValueError as exc:
        raise RuntimeError(f"{provider.name} returned non-JSON response.") from exc

    if not isinstance(data, dict):
        raise RuntimeError(f"{provider.name} returned malformed JSON response.")
    return data


async def test_model_profile(profile: ModelProfile) -> Dict[str, Any]:
    provider = provider_from_profile(profile)
    payload = {
        "model": profile.model,
        "messages": [{"role": "user", "content": "ping"}],
        "stream": False,
        "temperature": 0,
        "max_tokens": 8,
    }
    started = time.perf_counter()
    try:
        data = await call_upstream(provider, payload)
        text = extract_text_from_response(data)
        elapsed_ms = int((time.perf_counter() - started) * 1000)
        return {
            "ok": True,
            "elapsed_ms": elapsed_ms,
            "profile_id": profile.id,
            "model": profile.model,
            "message": text[:200],
        }
    except Exception as exc:
        elapsed_ms = int((time.perf_counter() - started) * 1000)
        return {
            "ok": False,
            "elapsed_ms": elapsed_ms,
            "profile_id": profile.id,
            "model": profile.model,
            "error": str(exc)[:1000],
        }


mcp = FastMCP(
    SERVER_NAME,
    instructions=(
        "Local BYOK multimodal chat gateway. Use smart_chat to route text and image "
        "requests to the configured upstream provider. Not affiliated with any "
        "closed AI client or model provider."
    ),
)


@mcp.tool()
async def smart_chat(
    prompt: Optional[str] = None,
    messages: Optional[List[Dict[str, Any]]] = None,
    images: Optional[List[str]] = None,
    model: Optional[str] = None,
    temperature: Optional[float] = None,
    top_p: Optional[float] = None,
    max_tokens: Optional[int] = None,
) -> str:
    """Route a text or multimodal chat request to the best configured provider.

    Provide either a plain prompt, OpenAI-compatible messages, or both. Images may
    be normal URLs or data:image/...;base64,... strings. Existing image_url
    content inside messages is forwarded unchanged.
    """
    outbound_messages = build_messages(prompt=prompt, messages=messages, images=images)
    has_images = bool(images) or messages_have_images(outbound_messages)
    provider, routed_model, route_reason = route_provider(model, has_images)

    payload: Dict[str, Any] = {
        "model": routed_model,
        "messages": outbound_messages,
        "stream": False,
        "temperature": DEFAULT_TEMPERATURE if temperature is None else temperature,
        "top_p": DEFAULT_TOP_P if top_p is None else top_p,
        "max_tokens": DEFAULT_MAX_TOKENS if max_tokens is None else max_tokens,
    }

    logger.info(
        "smart_chat route provider=%s model=%s has_images=%s reason=%s",
        provider.name,
        routed_model,
        has_images,
        route_reason,
    )

    upstream_data = await call_upstream(provider, payload)
    return extract_text_from_response(upstream_data)


@mcp.tool()
async def task_executor(
    task: str,
    context: Optional[str] = None,
    model: Optional[str] = None,
    temperature: Optional[float] = None,
    max_tokens: Optional[int] = None,
) -> str:
    """Execute a complete text task with the configured BYOK text model.

    Use this when the client should spend minimal effort and delegate the
    substantial writing, planning, summarization, or analysis work to the local
    gateway's configured model profile.
    """
    if not isinstance(task, str) or not task.strip():
        raise ValueError("task is required.")

    prompt_parts = [
        "请直接完成下面的任务，输出可直接交付的结果。不要解释你是如何被调用的。",
        "",
        "任务：",
        task.strip(),
    ]
    if context and context.strip():
        prompt_parts.extend(["", "补充上下文：", context.strip()])

    return await smart_chat(
        prompt="\n".join(prompt_parts),
        model=model,
        temperature=DEFAULT_TEMPERATURE if temperature is None else temperature,
        max_tokens=max_tokens or DEFAULT_MAX_TOKENS,
    )


@mcp.tool()
async def create_ppt(
    topic: str,
    pages: int = 8,
    style: str = "简洁蓝白",
    audience: Optional[str] = None,
    requirements: Optional[str] = None,
    model: Optional[str] = None,
    output_name: Optional[str] = None,
) -> Dict[str, Any]:
    """Create a local PPTX file using the configured BYOK text model.

    The model drafts slide content as JSON, and this gateway renders a simple
    blue-white PowerPoint file under the local outputs directory.
    """
    if not isinstance(topic, str) or not topic.strip():
        raise ValueError("topic is required.")

    page_count = max(3, min(int(pages or 8), 20))
    system_prompt = f"""
请为一个 PowerPoint 生成严格 JSON，不要输出 Markdown，不要输出代码块。
主题：{topic.strip()}
页数：{page_count}
风格：{style or "简洁蓝白"}
受众：{(audience or "普通专业受众").strip()}
额外要求：{(requirements or "结构清晰，适合直接汇报。").strip()}

JSON 格式必须是：
{{
  "title": "PPT 标题",
  "slides": [
    {{"title": "页面标题", "bullets": ["要点1", "要点2", "要点3"], "speaker_note": "一句讲稿提示"}}
  ]
}}

要求：
- slides 数量必须是 {page_count - 1}，封面由程序自动生成。
- 每页 3 到 5 个 bullets。
- bullets 要具体，不要空泛口号。
""".strip()

    draft = await smart_chat(
        prompt=system_prompt,
        model=model,
        temperature=0.4,
        max_tokens=max(1400, page_count * 420),
    )
    data = extract_json_object(draft)
    title = str(data.get("title") or topic).strip()
    raw_slides = data.get("slides")
    if not isinstance(raw_slides, list) or not raw_slides:
        raise ValueError("Model JSON did not contain slides.")

    slides: List[Dict[str, Any]] = []
    for item in raw_slides[: page_count - 1]:
        if isinstance(item, dict):
            slides.append(item)
        else:
            slides.append({"title": str(item), "bullets": []})

    pptx_path = build_pptx_file(title=title, slides=slides, output_name=output_name)
    return {
        "ok": True,
        "title": title,
        "pages": len(slides) + 1,
        "path": str(pptx_path),
        "message": f"PPT 已生成：{pptx_path}",
    }


@mcp.tool()
async def toolbox(
    action: str,
    prompt: Optional[str] = None,
    task: Optional[str] = None,
    topic: Optional[str] = None,
    context: Optional[str] = None,
    images: Optional[List[str]] = None,
    model: Optional[str] = None,
    pages: int = 8,
    style: str = "简洁蓝白",
    audience: Optional[str] = None,
    requirements: Optional[str] = None,
    output_name: Optional[str] = None,
    temperature: Optional[float] = None,
    max_tokens: Optional[int] = None,
) -> Any:
    """Unified MCP toolbox entrypoint.

    action can be:
    - chat: route prompt/images to the configured model.
    - task: delegate a complete task to the configured text model.
    - ppt: generate a local PPTX file.
    """
    normalized = (action or "").strip().lower()
    if normalized in {"chat", "smart_chat"}:
        return await smart_chat(
            prompt=prompt or task or topic,
            images=images,
            model=model,
            temperature=temperature,
            max_tokens=max_tokens,
        )
    if normalized in {"task", "execute", "task_executor"}:
        task_text = task or prompt or topic
        if not task_text:
            raise ValueError("task, prompt, or topic is required for action=task.")
        return await task_executor(
            task=task_text,
            context=context,
            model=model,
            temperature=temperature,
            max_tokens=max_tokens,
        )
    if normalized in {"ppt", "create_ppt", "presentation"}:
        ppt_topic = topic or task or prompt
        if not ppt_topic:
            raise ValueError("topic, task, or prompt is required for action=ppt.")
        return await create_ppt(
            topic=ppt_topic,
            pages=pages,
            style=style,
            audience=audience,
            requirements=requirements or context,
            model=model,
            output_name=output_name,
        )
    raise ValueError("Unsupported action. Use chat, task, or ppt.")


@mcp.tool()
async def list_skills() -> List[Dict[str, str]]:
    """List built-in local skills bundled with this MCP server."""
    return list_local_skills()


@mcp.tool()
async def run_skill(
    skill: str,
    task: str,
    context: Optional[str] = None,
    model: Optional[str] = None,
    pages: int = 8,
    output_name: Optional[str] = None,
    max_tokens: Optional[int] = None,
) -> Any:
    """Run a built-in local skill with the configured BYOK model.

    Built-in skills are lightweight instruction packs stored under skills/*.
    """
    if not isinstance(skill, str) or not skill.strip():
        raise ValueError("skill is required.")
    if not isinstance(task, str) or not task.strip():
        raise ValueError("task is required.")

    skill_data = read_skill(skill)
    skill_name = skill_data["name"]

    if skill_name == "ppt_writer":
        return await create_ppt(
            topic=task,
            pages=pages,
            style="简洁蓝白",
            requirements=(skill_data["instructions"] + "\n\n" + (context or "")).strip(),
            model=model,
            output_name=output_name,
        )

    prompt = f"""
你正在执行一个内置 MCP skill。请严格遵守 skill 说明，并直接输出可交付结果。

<skill name="{skill_name}">
{skill_data["instructions"]}
</skill>

<task>
{task.strip()}
</task>

<context>
{(context or "").strip()}
</context>
""".strip()

    return await smart_chat(
        prompt=prompt,
        model=model,
        temperature=0.3,
        max_tokens=max_tokens or DEFAULT_MAX_TOKENS,
    )


def detect_skill_route(task: str, context: Optional[str] = None) -> Tuple[str, str]:
    text = f"{task or ''}\n{context or ''}".lower()
    ppt_tokens = ("ppt", "powerpoint", "幻灯片", "演示文稿", "路演", "汇报稿", "slide", "deck")
    paper_tokens = ("论文", "paper", "摘要", "引言", "文献综述", "方法论", "实验分析", "学术", "润色")
    review_tokens = ("代码审查", "code review", "review code", "审查代码", "找 bug", "找bug", "安全风险", "缺少测试")
    plan_tokens = ("计划", "规划", "拆解", "步骤", "roadmap", "milestone", "待办", "任务分解")

    if any(token in text for token in ppt_tokens):
        return "ppt_writer", "matched presentation/PPT keywords"
    if any(token in text for token in paper_tokens):
        return "paper_writer", "matched paper/writing keywords"
    if any(token in text for token in review_tokens):
        return "code_reviewer", "matched code review keywords"
    if any(token in text for token in plan_tokens):
        return "task_planner", "matched planning keywords"
    return "task_executor", "no skill keyword matched; using generic task executor"


@mcp.tool()
async def auto_skill(
    task: str,
    context: Optional[str] = None,
    model: Optional[str] = None,
    pages: int = 8,
    output_name: Optional[str] = None,
    max_tokens: Optional[int] = None,
) -> Dict[str, Any]:
    """Automatically route a task to the best built-in skill.

    Use this as the default entrypoint when the user asks for PPT, paper
    writing, task planning, code review, or long-form work and does not want to
    choose a specific skill manually.
    """
    if not isinstance(task, str) or not task.strip():
        raise ValueError("task is required.")

    route, reason = detect_skill_route(task, context)
    if route == "task_executor":
        result = await task_executor(
            task=task,
            context=context,
            model=model,
            max_tokens=max_tokens,
        )
    else:
        result = await run_skill(
            skill=route,
            task=task,
            context=context,
            model=model,
            pages=pages,
            output_name=output_name,
            max_tokens=max_tokens,
        )
    return {
        "ok": True,
        "route": route,
        "reason": reason,
        "result": result,
    }


app = FastAPI(
    title="Smart Multimodal MCP SSE Gateway",
    description=(
        "A local BYOK MCP SSE gateway for routing text and multimodal chat requests. "
        "Not affiliated with any model provider or client vendor."
    ),
    version="0.1.0",
)


SETTINGS_HTML = """
<!doctype html>
<html lang="zh-CN">
<head>
  <meta charset="utf-8" />
  <meta name="viewport" content="width=device-width, initial-scale=1" />
  <title>Smart Multimodal MCP Gateway</title>
  <style>
    :root {
      color-scheme: light;
      --bg: #f3f8ff;
      --panel: #ffffff;
      --text: #102a43;
      --muted: #5f7189;
      --line: #c9ddf5;
      --accent: #1476d4;
      --accent-dark: #0c5da8;
      --ok: #0f7a52;
      --warn: #b76b00;
      --danger: #b42318;
      --shadow: 0 10px 30px rgba(20, 118, 212, 0.10);
    }

    * { box-sizing: border-box; }

    body {
      margin: 0;
      min-height: 100vh;
      background:
        linear-gradient(180deg, #eaf4ff 0, var(--bg) 260px),
        var(--bg);
      color: var(--text);
      font-family: "Segoe UI", system-ui, -apple-system, BlinkMacSystemFont, sans-serif;
      letter-spacing: 0;
    }

    header {
      border-bottom: 1px solid var(--line);
      background: rgba(255, 255, 255, 0.92);
      backdrop-filter: blur(12px);
      position: sticky;
      top: 0;
      z-index: 5;
    }

    .wrap {
      width: min(1180px, calc(100vw - 32px));
      margin: 0 auto;
    }

    .topbar {
      min-height: 72px;
      display: flex;
      align-items: center;
      justify-content: space-between;
      gap: 18px;
    }

    .title h1 {
      margin: 0;
      font-size: 22px;
      line-height: 1.2;
      font-weight: 700;
    }

    .title p {
      margin: 6px 0 0;
      color: var(--muted);
      font-size: 13px;
    }

    .status-row {
      display: flex;
      align-items: center;
      justify-content: flex-end;
      gap: 10px;
      flex-wrap: wrap;
    }

    .pill {
      display: inline-flex;
      align-items: center;
      min-height: 30px;
      padding: 0 10px;
      border-radius: 999px;
      border: 1px solid var(--line);
      background: #fff;
      color: var(--muted);
      font-size: 12px;
      white-space: nowrap;
    }

    .pill.ok { color: var(--ok); border-color: rgba(21, 115, 71, 0.28); }
    .pill.warn { color: var(--warn); border-color: rgba(181, 71, 8, 0.28); }

    main {
      padding: 28px 0 44px;
    }

    form {
      display: grid;
      gap: 18px;
    }

    .grid {
      display: grid;
      grid-template-columns: repeat(3, minmax(0, 1fr));
      gap: 18px;
      align-items: start;
    }

    section {
      background: var(--panel);
      border: 1px solid var(--line);
      border-radius: 8px;
      box-shadow: var(--shadow);
      padding: 18px;
    }

    section h2 {
      margin: 0;
      font-size: 16px;
      line-height: 1.3;
    }

    section .hint {
      margin: 6px 0 16px;
      color: var(--muted);
      font-size: 12px;
      line-height: 1.5;
    }

    label {
      display: grid;
      gap: 6px;
      margin: 12px 0;
      font-size: 12px;
      font-weight: 600;
      color: #344054;
    }

    input {
      width: 100%;
      min-height: 38px;
      border: 1px solid var(--line);
      border-radius: 6px;
      padding: 8px 10px;
      font: inherit;
      font-size: 13px;
      color: var(--text);
      background: #fff;
    }

    input:focus {
      outline: 2px solid rgba(15, 108, 189, 0.18);
      border-color: var(--accent);
    }

    .full {
      grid-column: 1 / -1;
    }

    .split {
      display: grid;
      grid-template-columns: repeat(3, minmax(0, 1fr));
      gap: 14px;
    }

    .actions {
      position: sticky;
      bottom: 0;
      z-index: 4;
      margin-top: 4px;
      padding: 14px 0;
      background: linear-gradient(180deg, rgba(246, 247, 249, 0), var(--bg) 24%);
    }

    .action-inner {
      display: flex;
      align-items: center;
      justify-content: space-between;
      gap: 14px;
      background: var(--panel);
      border: 1px solid var(--line);
      border-radius: 8px;
      box-shadow: var(--shadow);
      padding: 12px;
    }

    .message {
      min-height: 20px;
      font-size: 13px;
      color: var(--muted);
      overflow-wrap: anywhere;
    }

    .message.ok { color: var(--ok); }
    .message.error { color: var(--danger); }

    button {
      min-height: 38px;
      border: 0;
      border-radius: 6px;
      padding: 0 16px;
      background: var(--accent);
      color: #fff;
      font: inherit;
      font-size: 13px;
      font-weight: 700;
      cursor: pointer;
      white-space: nowrap;
    }

    button:hover { background: var(--accent-dark); }
    button.secondary {
      background: #fff;
      color: var(--text);
      border: 1px solid var(--line);
    }
    button.secondary:hover { background: #f2f4f7; }

    .button-row {
      display: flex;
      align-items: center;
      gap: 10px;
      flex-wrap: wrap;
      justify-content: flex-end;
    }

    @media (max-width: 920px) {
      .grid, .split { grid-template-columns: 1fr; }
      .topbar { align-items: flex-start; flex-direction: column; padding: 16px 0; }
      .status-row { justify-content: flex-start; }
      .action-inner { align-items: stretch; flex-direction: column; }
      .button-row { justify-content: flex-start; }
    }
  </style>
</head>
<body>
  <header>
    <div class="wrap topbar">
      <div class="title">
        <h1>Smart Multimodal MCP Gateway</h1>
        <p>本地 BYOK 配置面板。保存后 provider 路由立即读取新配置，端口和服务名需要重启。</p>
      </div>
      <div class="status-row">
        <span id="deepseekStatus" class="pill" data-label="文本模型">文本模型</span>
        <span id="dashscopeStatus" class="pill" data-label="多模态模型">多模态模型</span>
        <span id="openaiStatus" class="pill" data-label="备用多模态">备用多模态</span>
      </div>
    </div>
  </header>

  <main class="wrap">
    <form id="settingsForm">
      <div class="grid">
        <section>
          <h2>文本模型</h2>
          <p class="hint">纯文本默认渠道，适合日常聊天和代码任务。</p>
          <label>API Key <input name="DEEPSEEK_API_KEY" type="password" autocomplete="off" /></label>
          <label>Base URL <input name="DEEPSEEK_BASE_URL" /></label>
          <label>默认文本模型 <input name="DEEPSEEK_TEXT_MODEL" /></label>
        </section>

        <section>
          <h2>多模态模型</h2>
          <p class="hint">默认多模态渠道，带图片时优先使用。</p>
          <label>API Key <input name="DASHSCOPE_API_KEY" type="password" autocomplete="off" /></label>
          <label>Base URL <input name="DASHSCOPE_BASE_URL" /></label>
          <label>默认视觉模型 <input name="DASHSCOPE_VISION_MODEL" /></label>
        </section>

        <section>
          <h2>备用多模态</h2>
          <p class="hint">可配置另一个兼容接口，作为视觉模型或中转备用。</p>
          <label>API Key <input name="OPENAI_API_KEY" type="password" autocomplete="off" /></label>
          <label>Base URL <input name="OPENAI_BASE_URL" /></label>
          <label>默认视觉模型 <input name="OPENAI_VISION_MODEL" /></label>
        </section>
      </div>

      <section class="full">
        <h2>网关与连接池</h2>
        <p class="hint">大图片请求建议保持较长读写超时。HOST、PORT、服务名保存后需要重启服务才生效。</p>
        <div class="split">
          <label>服务名 <input name="MCP_SERVER_NAME" /></label>
          <label>HOST <input name="HOST" /></label>
          <label>PORT <input name="PORT" /></label>
          <label>Temperature <input name="DEFAULT_TEMPERATURE" /></label>
          <label>Top P <input name="DEFAULT_TOP_P" /></label>
          <label>Max Tokens <input name="DEFAULT_MAX_TOKENS" /></label>
          <label>连接超时 <input name="HTTP_CONNECT_TIMEOUT" /></label>
          <label>读取超时 <input name="HTTP_READ_TIMEOUT" /></label>
          <label>写入超时 <input name="HTTP_WRITE_TIMEOUT" /></label>
          <label>连接池等待 <input name="HTTP_POOL_TIMEOUT" /></label>
          <label>最大连接数 <input name="HTTP_MAX_CONNECTIONS" /></label>
          <label>Keep-Alive 数 <input name="HTTP_MAX_KEEPALIVE_CONNECTIONS" /></label>
        </div>
      </section>

      <div class="actions">
        <div class="action-inner">
          <div id="message" class="message">配置文件保存在本服务目录的 .env。</div>
          <div class="button-row">
            <button type="button" class="secondary" id="reloadButton">重新读取</button>
            <button type="submit">保存配置</button>
          </div>
        </div>
      </div>
    </form>
  </main>

  <script>
    const form = document.getElementById("settingsForm");
    const message = document.getElementById("message");
    const reloadButton = document.getElementById("reloadButton");

    function setMessage(text, kind = "") {
      message.textContent = text;
      message.className = "message" + (kind ? " " + kind : "");
    }

    function setProviderStatus(id, configured) {
      const el = document.getElementById(id);
      el.classList.toggle("ok", Boolean(configured));
      el.classList.toggle("warn", !configured);
      const name = el.dataset.label || el.textContent;
      el.textContent = name + (configured ? " 已配置" : " 未配置");
    }

    function fillForm(values) {
      for (const [key, value] of Object.entries(values)) {
        const input = form.elements[key];
        if (input) {
          input.value = value ?? "";
          if (value === "__KEEP_SECRET__") input.placeholder = "已配置，留空或保持不变将继续保留";
        }
      }
    }

    function collectForm() {
      const data = {};
      for (const element of form.elements) {
        if (element.name) data[element.name] = element.value;
      }
      return data;
    }

    async function loadConfig() {
      const res = await fetch("/config");
      if (!res.ok) throw new Error("读取配置失败");
      const data = await res.json();
      fillForm(data.values);
      setProviderStatus("deepseekStatus", data.configured.deepseek);
      setProviderStatus("dashscopeStatus", data.configured.dashscope);
      setProviderStatus("openaiStatus", data.configured.openai);
      setMessage(data.env_exists ? "已读取 .env。" : "当前没有 .env，保存后会创建。");
    }

    form.addEventListener("submit", async (event) => {
      event.preventDefault();
      setMessage("正在保存...");
      try {
        const res = await fetch("/config", {
          method: "POST",
          headers: { "Content-Type": "application/json" },
          body: JSON.stringify(collectForm()),
        });
        const data = await res.json();
        if (!res.ok) throw new Error(data.detail || "保存失败");
        fillForm(data.values);
        setProviderStatus("deepseekStatus", data.configured.deepseek);
        setProviderStatus("dashscopeStatus", data.configured.dashscope);
        setProviderStatus("openaiStatus", data.configured.openai);
        setMessage("保存成功。Provider 配置已应用；HOST、PORT、服务名需要重启服务。", "ok");
      } catch (error) {
        setMessage(error.message || String(error), "error");
      }
    });

    reloadButton.addEventListener("click", () => {
      loadConfig().catch((error) => setMessage(error.message || String(error), "error"));
    });

    loadConfig().catch((error) => setMessage(error.message || String(error), "error"));
  </script>
</body>
</html>
"""

SETTINGS_HTML = """
<!doctype html>
<html lang="zh-CN">
<head>
  <meta charset="utf-8" />
  <meta name="viewport" content="width=device-width, initial-scale=1" />
  <title>模型档案 - Smart MCP Gateway</title>
  <style>
    :root {
      --bg: #f6f7f9;
      --panel: #ffffff;
      --text: #172033;
      --muted: #667085;
      --line: #d9dee8;
      --accent: #0f6cbd;
      --accent-dark: #0b5aa0;
      --ok: #157347;
      --warn: #b54708;
      --danger: #b42318;
      --shadow: 0 10px 30px rgba(23, 32, 51, 0.08);
    }
    * { box-sizing: border-box; }
    body {
      margin: 0;
      min-height: 100vh;
      background: var(--bg);
      color: var(--text);
      font-family: "Segoe UI", system-ui, -apple-system, BlinkMacSystemFont, sans-serif;
      letter-spacing: 0;
    }
    header {
      position: sticky;
      top: 0;
      z-index: 10;
      border-bottom: 1px solid var(--line);
      background: rgba(255,255,255,.94);
      backdrop-filter: blur(12px);
    }
    .wrap { width: min(1180px, calc(100vw - 32px)); margin: 0 auto; }
    .topbar {
      min-height: 72px;
      display: flex;
      justify-content: space-between;
      align-items: center;
      gap: 18px;
    }
    h1 { margin: 0; font-size: 22px; line-height: 1.2; }
    .subtitle { margin: 6px 0 0; color: var(--muted); font-size: 13px; }
    main { padding: 28px 0 46px; }
    .status-row { display: flex; gap: 10px; flex-wrap: wrap; justify-content: flex-end; }
    .pill {
      min-height: 30px;
      display: inline-flex;
      align-items: center;
      padding: 0 10px;
      border: 1px solid var(--line);
      border-radius: 999px;
      background: #fff;
      color: var(--muted);
      font-size: 12px;
      white-space: nowrap;
    }
    .pill.ok { color: var(--ok); border-color: rgba(21,115,71,.28); }
    .pill.warn { color: var(--warn); border-color: rgba(181,71,8,.28); }
    .layout { display: grid; grid-template-columns: 280px minmax(0, 1fr); gap: 18px; align-items: start; }
    section, .profile-card {
      background: var(--panel);
      border: 1px solid var(--line);
      border-radius: 8px;
      box-shadow: var(--shadow);
      padding: 16px;
    }
    section h2 { margin: 0 0 8px; font-size: 16px; }
    .hint { color: var(--muted); font-size: 12px; line-height: 1.5; margin: 0 0 14px; }
    label { display: grid; gap: 6px; margin: 12px 0; color: #344054; font-size: 12px; font-weight: 600; }
    input, select {
      width: 100%;
      min-height: 38px;
      border: 1px solid var(--line);
      border-radius: 6px;
      padding: 8px 10px;
      background: #fff;
      color: var(--text);
      font: inherit;
      font-size: 13px;
    }
    input:focus, select:focus { outline: 2px solid rgba(15,108,189,.18); border-color: var(--accent); }
    .profiles { display: grid; gap: 14px; }
    .profile-card header {
      position: static;
      border: 0;
      background: transparent;
      display: flex;
      justify-content: space-between;
      align-items: flex-start;
      gap: 12px;
      margin-bottom: 8px;
    }
    .profile-title { margin: 0; font-size: 16px; }
    .profile-meta { margin: 4px 0 0; color: var(--muted); font-size: 12px; overflow-wrap: anywhere; }
    .form-grid { display: grid; grid-template-columns: repeat(2, minmax(0, 1fr)); gap: 0 14px; }
    .wide { grid-column: 1 / -1; }
    .model-switch {
      display: grid;
      grid-template-columns: minmax(0, 1fr) auto;
      gap: 8px;
      align-items: end;
    }
    .key-row { display: grid; grid-template-columns: minmax(0, 1fr) auto; gap: 8px; align-items: end; }
    .key-actions { display: flex; gap: 8px; flex-wrap: wrap; align-items: end; }
    .switch-row { display: flex; gap: 8px; align-items: center; margin: 12px 0; color: #344054; font-size: 12px; font-weight: 600; }
    .switch-row input { width: 18px; min-height: 18px; }
    .actions {
      position: sticky;
      bottom: 0;
      z-index: 9;
      margin-top: 18px;
      padding: 14px 0;
      background: linear-gradient(180deg, rgba(246,247,249,0), var(--bg) 28%);
    }
    .action-inner {
      display: flex;
      justify-content: space-between;
      align-items: center;
      gap: 12px;
      padding: 12px;
      border: 1px solid var(--line);
      border-radius: 8px;
      background: var(--panel);
      box-shadow: var(--shadow);
    }
    .buttons { display: flex; gap: 8px; flex-wrap: wrap; justify-content: flex-end; }
    button {
      min-height: 38px;
      border: 0;
      border-radius: 6px;
      padding: 0 14px;
      background: var(--accent);
      color: white;
      font: inherit;
      font-size: 13px;
      font-weight: 700;
      cursor: pointer;
      white-space: nowrap;
    }
    button:hover { background: var(--accent-dark); }
    button.secondary { background: #fff; color: var(--text); border: 1px solid var(--line); }
    button.secondary:hover { background: #f2f4f7; }
    button.danger { background: #fff; color: var(--danger); border: 1px solid rgba(180,35,24,.35); }
    button.danger:hover { background: #fff5f5; }
    button.small { min-height: 34px; padding: 0 10px; font-size: 12px; }
    .message { min-height: 20px; color: var(--muted); font-size: 13px; overflow-wrap: anywhere; }
    .message.ok { color: var(--ok); }
    .message.error { color: var(--danger); }
    .test-result { margin-top: 8px; font-size: 12px; color: var(--muted); overflow-wrap: anywhere; }
    .test-result.ok { color: var(--ok); }
    .test-result.error { color: var(--danger); }
    @media (max-width: 920px) {
      .topbar, .action-inner { align-items: stretch; flex-direction: column; }
      .status-row, .buttons { justify-content: flex-start; }
      .layout, .form-grid { grid-template-columns: 1fr; }
    }
  </style>
</head>
<body>
  <header>
    <div class="wrap topbar">
      <div>
        <h1>模型档案</h1>
        <p class="subtitle">管理自定义 API、模型切换和可用性测试。密钥默认遮盖，需要时可手动读取和复制。</p>
      </div>
      <div class="status-row">
        <span id="textStatus" class="pill">文本模型</span>
        <span id="multiStatus" class="pill">多模态模型</span>
        <span id="profileCount" class="pill">0 个档案</span>
      </div>
    </div>
  </header>

  <main class="wrap">
    <div class="layout">
      <section>
        <h2>默认路由</h2>
        <p class="hint">纯文本默认走文本模型；包含图片时默认走多模态模型。也可以通过 model 精确命中档案名或模型名。</p>
        <label>默认文本模型 <select id="defaultText"></select></label>
        <label>默认多模态模型 <select id="defaultMultimodal"></select></label>
        <button type="button" id="addProfile">新增模型档案</button>
      </section>

      <div>
        <div id="profiles" class="profiles"></div>
        <div class="actions">
          <div class="action-inner">
            <div id="message" class="message">配置保存在本地 models.json。</div>
            <div class="buttons">
              <button type="button" class="secondary" id="reloadModels">重新读取</button>
              <button type="button" id="saveModels">保存模型档案</button>
            </div>
          </div>
        </div>
      </div>
    </div>
  </main>

  <script>
    const KEEP_SECRET = "__KEEP_SECRET__";
    const profilesEl = document.getElementById("profiles");
    const defaultTextEl = document.getElementById("defaultText");
    const defaultMultiEl = document.getElementById("defaultMultimodal");
    const messageEl = document.getElementById("message");
    let state = { profiles: [], default_text_profile_id: "", default_multimodal_profile_id: "" };

    function uid() {
      return "model-" + Math.random().toString(36).slice(2, 10);
    }

    function setMessage(text, kind = "") {
      messageEl.textContent = text;
      messageEl.className = "message" + (kind ? " " + kind : "");
    }

    function capabilityName(value) {
      return { text: "文本模型", multimodal: "多模态模型", fallback: "备用模型" }[value] || value;
    }

    function newProfile() {
      return {
        id: uid(),
        display_name: "新模型",
        capability: "text",
        base_url: "",
        api_key: "",
        model: "",
        model_options: [],
        enabled: true,
      };
    }

    function profileTitle(profile) {
      return (profile.display_name || profile.model || profile.id || "未命名模型").trim();
    }

    function renderDefaults() {
      const enabled = state.profiles.filter((p) => p.enabled);
      const textOptions = enabled.filter((p) => p.capability === "text");
      const multiOptions = enabled.filter((p) => p.capability === "multimodal" || p.capability === "fallback");
      fillSelect(defaultTextEl, textOptions.length ? textOptions : enabled, state.default_text_profile_id);
      fillSelect(defaultMultiEl, multiOptions.length ? multiOptions : enabled, state.default_multimodal_profile_id);

      const textOk = Boolean(state.default_text_profile_id);
      const multiOk = Boolean(state.default_multimodal_profile_id);
      setPill("textStatus", "文本模型", textOk);
      setPill("multiStatus", "多模态模型", multiOk);
      document.getElementById("profileCount").textContent = state.profiles.length + " 个档案";
    }

    function setPill(id, label, ok) {
      const el = document.getElementById(id);
      el.classList.toggle("ok", ok);
      el.classList.toggle("warn", !ok);
      el.textContent = label + (ok ? " 已选择" : " 未选择");
    }

    function fillSelect(select, profiles, selected) {
      select.innerHTML = "";
      for (const profile of profiles) {
        const option = document.createElement("option");
        option.value = profile.id;
        option.textContent = profileTitle(profile) + " · " + (profile.model || "未填写模型名");
        option.selected = profile.id === selected;
        select.appendChild(option);
      }
    }

    function renderProfiles() {
      profilesEl.innerHTML = "";
      for (const profile of state.profiles) {
        profilesEl.appendChild(renderProfile(profile));
      }
      renderDefaults();
    }

    function renderProfile(profile) {
      const card = document.createElement("article");
      card.className = "profile-card";
      card.dataset.id = profile.id;
      card.innerHTML = `
        <header>
          <div>
            <h3 class="profile-title">${escapeHtml(profileTitle(profile))}</h3>
            <p class="profile-meta">${escapeHtml(capabilityName(profile.capability))} · ${escapeHtml(profile.model || "未填写模型名")}</p>
          </div>
          <div class="buttons">
            <button type="button" class="small secondary test-btn">测试模型</button>
            <button type="button" class="small danger delete-btn">删除</button>
          </div>
        </header>
        <div class="form-grid">
          <label>显示名 <input data-field="display_name" value="${escapeAttr(profile.display_name)}" /></label>
          <label>能力类型
            <select data-field="capability">
              <option value="text">文本模型</option>
              <option value="multimodal">多模态模型</option>
              <option value="fallback">备用模型</option>
            </select>
          </label>
          <label>Base URL <input data-field="base_url" value="${escapeAttr(profile.base_url)}" /></label>
          <label class="model-switch">
            <span>当前模型 <input data-field="model" value="${escapeAttr(profile.model)}" /></span>
            <button type="button" class="small secondary switch-model">切换</button>
          </label>
          <label class="wide">可选模型列表 <input data-field="model_options_text" value="${escapeAttr((profile.model_options || []).join(', '))}" placeholder="例如 deepseek-chat, deepseek-reasoner, flash, pro" /></label>
          <label class="key-row wide">
            <span>API Key <input data-field="api_key" type="password" autocomplete="off" value="${escapeAttr(profile.api_key)}" placeholder="已配置时保持占位即可" /></span>
            <span class="key-actions">
              <button type="button" class="small secondary toggle-key">显示</button>
              <button type="button" class="small secondary reveal-key">读取已保存</button>
              <button type="button" class="small secondary copy-key">复制</button>
            </span>
          </label>
          <label class="switch-row"><input data-field="enabled" type="checkbox" ${profile.enabled ? "checked" : ""} /> 启用此档案</label>
        </div>
        <div class="test-result"></div>
      `;
      card.querySelector('[data-field="capability"]').value = profile.capability;
      card.querySelectorAll("[data-field]").forEach((input) => {
        input.addEventListener("input", () => syncCard(card));
        input.addEventListener("change", () => syncCard(card));
      });
      card.querySelector(".delete-btn").addEventListener("click", () => {
        state.profiles = state.profiles.filter((item) => item.id !== profile.id);
        if (state.default_text_profile_id === profile.id) state.default_text_profile_id = "";
        if (state.default_multimodal_profile_id === profile.id) state.default_multimodal_profile_id = "";
        renderProfiles();
      });
      card.querySelector(".toggle-key").addEventListener("click", (event) => {
        const input = card.querySelector('[data-field="api_key"]');
        input.type = input.type === "password" ? "text" : "password";
        event.target.textContent = input.type === "password" ? "显示" : "隐藏";
      });
      card.querySelector(".reveal-key").addEventListener("click", () => revealSavedKey(card));
      card.querySelector(".copy-key").addEventListener("click", () => copyProfileKey(card));
      card.querySelector(".switch-model").addEventListener("click", () => switchModel(card));
      card.querySelector(".test-btn").addEventListener("click", () => testProfile(card));
      return card;
    }

    function syncCard(card) {
      const profile = state.profiles.find((item) => item.id === card.dataset.id);
      if (!profile) return;
      card.querySelectorAll("[data-field]").forEach((input) => {
        const field = input.dataset.field;
        if (field === "model_options_text") {
          profile.model_options = input.value.split(",").map((item) => item.trim()).filter(Boolean);
        } else {
          profile[field] = input.type === "checkbox" ? input.checked : input.value;
        }
      });
      if (profile.model && !profile.model_options.includes(profile.model)) {
        profile.model_options = [profile.model, ...profile.model_options];
        const optionsInput = card.querySelector('[data-field="model_options_text"]');
        if (optionsInput) optionsInput.value = profile.model_options.join(", ");
      }
      card.querySelector(".profile-title").textContent = profileTitle(profile);
      card.querySelector(".profile-meta").textContent = capabilityName(profile.capability) + " · " + (profile.model || "未填写模型名");
      renderDefaults();
    }

    async function revealSavedKey(card) {
      syncCard(card);
      const input = card.querySelector('[data-field="api_key"]');
      const profile = state.profiles.find((item) => item.id === card.dataset.id);
      if (!profile || !profile.id) {
        throw new Error("请先保存这个模型档案，再读取已保存 Key。");
      }
      const res = await fetch(`/models/${encodeURIComponent(profile.id)}/api-key`);
      const data = await res.json();
      if (!res.ok) throw new Error(data.detail || "读取 API Key 失败");
      input.value = data.api_key || "";
      input.type = "text";
      card.querySelector(".toggle-key").textContent = "隐藏";
      profile.api_key = input.value;
      return input.value;
    }

    async function copyProfileKey(card) {
      try {
        const input = card.querySelector('[data-field="api_key"]');
        let value = input.value;
        if (!value || value === KEEP_SECRET) value = await revealSavedKey(card);
        if (!value) throw new Error("当前档案没有可复制的 API Key。");
        await navigator.clipboard.writeText(value);
        card.querySelector(".test-result").textContent = "API Key 已复制。";
        card.querySelector(".test-result").className = "test-result ok";
      } catch (error) {
        card.querySelector(".test-result").textContent = error.message || String(error);
        card.querySelector(".test-result").className = "test-result error";
      }
    }

    function switchModel(card) {
      syncCard(card);
      const profile = state.profiles.find((item) => item.id === card.dataset.id);
      if (!profile) return;
      const options = (profile.model_options || []).filter(Boolean);
      if (!options.length) {
        card.querySelector(".test-result").textContent = "请先填写可选模型列表，例如 flash, pro。";
        card.querySelector(".test-result").className = "test-result error";
        return;
      }
      const currentIndex = Math.max(0, options.indexOf(profile.model));
      const next = options[(currentIndex + 1) % options.length];
      profile.model = next;
      card.querySelector('[data-field="model"]').value = next;
      card.querySelector(".profile-meta").textContent = capabilityName(profile.capability) + " · " + next;
      renderDefaults();
    }

    function collectState() {
      return {
        profiles: state.profiles,
        default_text_profile_id: defaultTextEl.value,
        default_multimodal_profile_id: defaultMultiEl.value,
      };
    }

    async function loadModels() {
      const res = await fetch("/models");
      if (!res.ok) throw new Error("读取模型档案失败");
      const data = await res.json();
      state = data.settings;
      renderProfiles();
      setMessage(data.models_exists ? "已读取 models.json。" : "未发现 models.json，正在使用旧配置生成的默认档案。");
    }

    async function saveModels() {
      setMessage("正在保存...");
      const res = await fetch("/models", {
        method: "POST",
        headers: { "Content-Type": "application/json" },
        body: JSON.stringify(collectState()),
      });
      const data = await res.json();
      if (!res.ok) throw new Error(data.detail || "保存失败");
      state = data.settings;
      renderProfiles();
      setMessage("保存成功，路由已立即生效。", "ok");
    }

    async function testProfile(card) {
      syncCard(card);
      const resultEl = card.querySelector(".test-result");
      const profile = state.profiles.find((item) => item.id === card.dataset.id);
      resultEl.textContent = "测试中...";
      resultEl.className = "test-result";
      const res = await fetch("/models/test", {
        method: "POST",
        headers: { "Content-Type": "application/json" },
        body: JSON.stringify({ profile }),
      });
      const data = await res.json();
      if (data.ok) {
        resultEl.textContent = `可用 · ${data.elapsed_ms}ms · ${data.message || "ok"}`;
        resultEl.className = "test-result ok";
      } else {
        resultEl.textContent = `不可用 · ${data.elapsed_ms || 0}ms · ${data.error || data.detail || "测试失败"}`;
        resultEl.className = "test-result error";
      }
    }

    function escapeHtml(value) {
      return String(value ?? "").replace(/[&<>"']/g, (ch) => ({ "&": "&amp;", "<": "&lt;", ">": "&gt;", '"': "&quot;", "'": "&#39;" }[ch]));
    }
    function escapeAttr(value) { return escapeHtml(value); }

    document.getElementById("addProfile").addEventListener("click", () => {
      state.profiles.push(newProfile());
      renderProfiles();
    });
    document.getElementById("reloadModels").addEventListener("click", () => {
      loadModels().catch((error) => setMessage(error.message || String(error), "error"));
    });
    document.getElementById("saveModels").addEventListener("click", () => {
      saveModels().catch((error) => setMessage(error.message || String(error), "error"));
    });
    defaultTextEl.addEventListener("change", () => state.default_text_profile_id = defaultTextEl.value);
    defaultMultiEl.addEventListener("change", () => state.default_multimodal_profile_id = defaultMultiEl.value);
    loadModels().catch((error) => setMessage(error.message || String(error), "error"));
  </script>
</body>
</html>
"""


@app.get("/", response_class=HTMLResponse)
async def settings_page() -> str:
    return SETTINGS_HTML


@app.get("/config")
async def get_config() -> Dict[str, Any]:
    return public_config()


@app.post("/config")
async def save_config(config: GatewayConfig) -> Dict[str, Any]:
    try:
        config = merge_secret_placeholders(config)
        write_env_file(config)
        apply_config_to_process(config)
        await close_http_client()
    except OSError as exc:
        raise HTTPException(status_code=500, detail=f"Failed to save .env: {exc}") from exc
    return public_config()


@app.get("/models")
async def get_models() -> Dict[str, Any]:
    return public_model_settings()


@app.get("/models/{profile_id}/api-key")
async def get_model_api_key(profile_id: str) -> Dict[str, str]:
    settings = read_model_settings(include_secrets=True)
    profile = find_profile_by_id(settings, profile_id)
    if profile is None:
        raise HTTPException(status_code=404, detail="Model profile not found.")
    return {"api_key": profile.api_key or ""}


@app.post("/models")
async def save_models(settings: ModelSettings) -> Dict[str, Any]:
    try:
        saved = write_model_settings(settings)
    except OSError as exc:
        raise HTTPException(status_code=500, detail=f"Failed to save models.json: {exc}") from exc
    return {
        "models_path": str(MODELS_PATH),
        "models_exists": MODELS_PATH.exists(),
        "settings": mask_model_settings(saved).model_dump(),
        "configured": {profile.id: bool(profile.api_key) for profile in saved.profiles},
    }


@app.post("/models/test")
async def test_model(request: ModelTestRequest) -> Dict[str, Any]:
    profile: Optional[ModelProfile] = None
    settings = read_model_settings(include_secrets=True)

    if request.profile is not None:
        old = find_profile_by_id(settings, request.profile.id)
        data = request.profile.model_dump()
        if data["api_key"] in {"", KEEP_SECRET} and old:
            data["api_key"] = old.api_key
        profile = ModelProfile(**data)
    elif request.profile_id:
        profile = find_profile_by_id(settings, request.profile_id)

    if profile is None:
        raise HTTPException(status_code=404, detail="Model profile not found.")
    if not profile.enabled:
        raise HTTPException(status_code=400, detail="Model profile is disabled.")
    if not profile.base_url or not profile.model:
        raise HTTPException(status_code=400, detail="Base URL and model are required.")

    return await test_model_profile(profile)


@app.get("/health")
async def health() -> Dict[str, Any]:
    model_settings = read_model_settings(include_secrets=True)
    return {
        "status": "ok",
        "server": SERVER_NAME,
        "transport": "mcp-sse",
        "sse_url": "/sse",
        "model_profiles": {
            "count": len(model_settings.profiles),
            "enabled": len([profile for profile in model_settings.profiles if profile.enabled]),
            "default_text_profile_id": model_settings.default_text_profile_id,
            "default_multimodal_profile_id": model_settings.default_multimodal_profile_id,
        },
        "providers": {
            "deepseek": {
                "configured": bool(env_str("DEEPSEEK_API_KEY")),
                "base_url": env_str("DEEPSEEK_BASE_URL", "https://api.deepseek.com/v1"),
                "default_model": env_str("DEEPSEEK_TEXT_MODEL", "deepseek-chat"),
            },
            "dashscope": {
                "configured": bool(env_str("DASHSCOPE_API_KEY")),
                "base_url": env_str(
                    "DASHSCOPE_BASE_URL",
                    "https://dashscope.aliyuncs.com/compatible-mode/v1",
                ),
                "default_model": env_str("DASHSCOPE_VISION_MODEL", "qwen-vl-max"),
            },
            "openai": {
                "configured": bool(env_str("OPENAI_API_KEY")),
                "base_url": env_str("OPENAI_BASE_URL", "https://api.openai.com/v1"),
                "default_model": env_str("OPENAI_VISION_MODEL", "gpt-4o"),
            },
        },
    }


@app.on_event("shutdown")
async def shutdown_event() -> None:
    await close_http_client()


app.mount("/", mcp.sse_app())


if __name__ == "__main__":
    uvicorn.run(
        "main:app",
        host=env_str("HOST", "0.0.0.0"),
        port=env_int("PORT", 8010),
        reload=False,
    )
